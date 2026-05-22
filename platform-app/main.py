import os
import json
from datetime import datetime
from typing import Optional, List

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
from fastapi.responses import JSONResponse, Response
from starlette.middleware.base import BaseHTTPMiddleware
from pydantic import BaseModel
import requests
from urllib.parse import urlparse
import re
import threading

# 导入数据库操作封装模块
import wxcloud_db

app = FastAPI(title="海外游戏平台运营策略中心")

# CORS中间件
from fastapi.middleware.cors import CORSMiddleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 增大请求体大小限制
MAX_UPLOAD_SIZE = 10 * 1024 * 1024 * 1024  # 已取消限制，设为10GB

class LimitUploadSizeMiddleware(BaseHTTPMiddleware):
    async def dispatch(self, request: Request, call_next):
        content_length = request.headers.get("content-length")
        if content_length and int(content_length) > MAX_UPLOAD_SIZE:
            return JSONResponse(
                status_code=413,
                content={"code": -1, "msg": f"文件过大，最大支持{MAX_UPLOAD_SIZE // (1024*1024)}MB"}
            )
        return await call_next(request)

app.add_middleware(LimitUploadSizeMiddleware)

# 临时缓存：存储解析报告时的原始内容
import uuid
_report_temp_cache = {}
_CACHE_MAX_AGE_SECONDS = 1800

def _cleanup_expired_cache():
    now = datetime.now()
    expired = [k for k, v in _report_temp_cache.items()
               if (now - v["created_at"]).total_seconds() > _CACHE_MAX_AGE_SECONDS]
    for k in expired:
        del _report_temp_cache[k]


# ==================== 数据库初始化 ====================

@app.on_event("startup")
async def startup_init():
    """启动时在后台线程执行策略数据初始化，不阻塞服务启动"""
    def _run_update():
        try:
            # 先记录当前数据库状态（用于对比）
            try:
                with wxcloud_db.get_db() as db:
                    cursor = db.cursor()
                    cursor.execute("SELECT COUNT(*) as cnt FROM strategies")
                    strat_cnt = cursor.fetchone()["cnt"]
                    cursor.execute("SELECT COUNT(*) as cnt FROM uploaded_docs")
                    doc_cnt = cursor.fetchone()["cnt"]
                    # 统计策略中的report_data总数
                    cursor.execute("SELECT id, content FROM strategies WHERE content IS NOT NULL AND content != ''")
                    rows = cursor.fetchall()
                    report_data_total = 0
                    for row in rows:
                        try:
                            content = json.loads(row["content"]) if isinstance(row["content"], str) else row["content"]
                            for sec in content.get("sections", []):
                                for item in sec.get("items", []):
                                    report_data_total += len(item.get("report_data", []))
                        except Exception:
                            pass
                    cursor.close()
                print(f"[startup] DB state BEFORE update: {strat_cnt} strategies, {doc_cnt} docs, {report_data_total} report_data entries")
            except Exception as e:
                print(f"[startup] Could not check DB state: {e}")

            import subprocess
            result = subprocess.run(
                ["python", "update_strategies.py"],
                capture_output=True, text=True, timeout=120
            )
            if result.returncode == 0:
                print("[startup] Strategies updated:", result.stdout.strip()[-300:])
            else:
                print("[startup] Strategy update failed:", result.stderr[-500:])

            # 更新后再次检查数据库状态
            try:
                with wxcloud_db.get_db() as db:
                    cursor = db.cursor()
                    cursor.execute("SELECT COUNT(*) as cnt FROM strategies")
                    strat_cnt_after = cursor.fetchone()["cnt"]
                    cursor.execute("SELECT id, content FROM strategies WHERE content IS NOT NULL AND content != ''")
                    rows = cursor.fetchall()
                    report_data_after = 0
                    for row in rows:
                        try:
                            content = json.loads(row["content"]) if isinstance(row["content"], str) else row["content"]
                            for sec in content.get("sections", []):
                                for item in sec.get("items", []):
                                    report_data_after += len(item.get("report_data", []))
                        except Exception:
                            pass
                    cursor.close()
                print(f"[startup] DB state AFTER update: {strat_cnt_after} strategies, {report_data_after} report_data entries")
                # 如果report_data数量减少了，发出警告
                if report_data_after < report_data_total:
                    print(f"[startup] ⚠️ WARNING: report_data decreased from {report_data_total} to {report_data_after}!")
            except Exception as e:
                print(f"[startup] Could not check DB state after update: {e}")

        except Exception as e:
            print(f"[startup] Init error: {e}")
    
    # 后台线程执行，不阻塞 HTTP 服务启动
    thread = threading.Thread(target=_run_update, daemon=True)
    thread.start()
    print("[startup] Strategy update started in background thread")


# ==================== 数据模型 ====================

class StrategyUpdate(BaseModel):
    platform_id: int
    category: str
    title: str
    content: str

class DocUpload(BaseModel):
    platform_id: Optional[int] = None
    filename: str
    content: str
    doc_type: str = "policy"
    game_name: Optional[str] = None

class UrlUpload(BaseModel):
    url: str
    platform_id: Optional[int] = None
    doc_type: str = "policy"
    title: Optional[str] = None
    game_name: Optional[str] = None


# ==================== 路由 ====================

@app.get("/healthz")
async def health_check():
    return {"status": "ok"}


@app.get("/api/platforms")
async def get_platforms():
    """获取所有平台列表"""
    try:
        platforms = wxcloud_db.get_all_platforms()
        return {"code": 0, "data": platforms}
    except Exception as e:
        return {"code": -1, "msg": f"加载平台数据失败: {str(e)}"}


def _normalize_sections(strategy):
    """统一将 section 中的 '标题' 字段映射为 'title'，兼容旧数据"""
    content = strategy.get("content")
    if not content:
        strategy["content"] = {"sections": []}
        return
    if isinstance(content, str):
        try:
            content = __import__("json").loads(content)
            strategy["content"] = content
        except Exception:
            strategy["content"] = {"sections": []}
            return
    sections = content.get("sections", [])
    for sec in sections:
        if "title" not in sec and "标题" in sec:
            sec["title"] = sec.pop("标题")


@app.get("/api/strategies")
async def get_all_strategies():
    """获取所有策略"""
    try:
        strategies = wxcloud_db.get_all_strategies()
        for s in strategies:
            _normalize_sections(s)
        return {"code": 0, "data": strategies}
    except Exception as e:
        return {"code": -1, "msg": f"加载策略失败: {str(e)}"}


@app.get("/api/strategies/{platform_id}")
async def get_strategies(platform_id: int):
    """获取指定平台的所有策略"""
    try:
        strategies = wxcloud_db.get_strategies_by_platform(platform_id)
        for s in strategies:
            _normalize_sections(s)
        return {"code": 0, "data": strategies}
    except Exception as e:
        return {"code": -1, "msg": f"加载策略失败: {str(e)}"}


# ==================== 报告解析相关函数 ====================

def parse_report_content(text: str, filename: str = "") -> dict:
    """解析报告内容，提炼关键数据和结论"""
    clean_text = text.replace('<!--TABLE_START-->', '').replace('<!--TABLE_END-->', '')

    # 预处理：去除幻灯片分页标记，合并碎片化短行
    clean_text = re.sub(r'---\s*第\s*\d+\s*页.*?---\s*\n?', '\n', clean_text)
    clean_text = re.sub(r'\[备注\]\s*', '', clean_text)

    # 合并碎片化短行：将连续的短行（<20字符且不含数字）合并为一行
    raw_lines = clean_text.strip().split('\n')
    merged_lines = []
    buffer = ""
    for line in raw_lines:
        stripped = line.strip()
        if not stripped:
            if buffer:
                merged_lines.append(buffer)
                buffer = ""
            continue
        # 如果当前行很短且不含关键数据，尝试与buffer合并
        has_number = bool(re.search(r'\d', stripped))
        if len(stripped) < 20 and not has_number and buffer and len(buffer) < 100:
            buffer = buffer + " " + stripped
        else:
            if buffer:
                merged_lines.append(buffer)
            buffer = stripped
    if buffer:
        merged_lines.append(buffer)

    lines = merged_lines
    key_data_points = []
    conclusions = []
    full_text = "\n".join(lines)
    full_lower = full_text.lower()

    wl_patterns = [
        r'(?:愿望单|wishlist)[^\n]*?[\d,.]+[^\n]*',
        r'[^\n]*?[\d,]+\s*(?:wishlists?|愿望单)[^\n]*',
        r'[^\n]*?got\s+[\d,]+\s+wishlist[^\n]*',
        r'获得愿望单[^\n]*?[\d,.]+[^\n]*',
    ]
    player_patterns = [
        r'(?:游玩人数|玩家数|players?|游玩数)[^\n]*?[\d,.]+[^\n]*',
        r'[^\n]*?[\d,]+\s*(?:players?|游玩人数|玩家)[^\n]*',
        r'(?:游玩人数|players?)[^\n]*',
    ]
    sales_patterns = [
        r'(?:下载量|销量|sales?|downloads?|copies)[^\n]*?[\d,.]+[^\n]*',
        r'[^\n]*?[\d,]+\s*(?:copies|downloads?|units?|销量)[^\n]*',
    ]
    rate_patterns = [
        r'(?:转化率|conversion|率)[^\n]*?[\d,.]+%[^\n]*',
        r'[^\n]*?[\d,.]+%[^\n]*?(?:转化|conversion|增长|growth|rate)[^\n]*',
    ]
    general_number_patterns = [
        r'[^\n]*?(?:增长|增加|提升|翻倍|doubl|triple|grew|increase|boost|surge)[^\n]*?[\d,.]+[^\n]*',
        r'[^\n]*?[\d,.]+[^\n]*?(?:增长|增加|提升|翻倍|doubl|triple|grew|increase|boost|surge)[^\n]*',
    ]

    all_patterns = wl_patterns + player_patterns + sales_patterns + rate_patterns + general_number_patterns

    seen_data = set()
    for pattern in all_patterns:
        matches = re.finditer(pattern, full_text, re.IGNORECASE)
        for m in matches:
            raw = m.group(0).strip()
            if len(raw) < 5 or len(raw) > 200:
                continue
            nums = re.findall(r'[\d,]+', raw)
            num_key = tuple(sorted(nums))
            if num_key and num_key not in seen_data:
                seen_data.add(num_key)
                key_data_points.append(raw.strip())

    achievement_keywords = [
        '榜单', '排名', 'top', 'ranking', 'chart', 'most played', 'most wished',
        '最受欢迎', '热门', 'popular', 'trending', 'featured', 'highlighted',
        '推荐', '首页', '展示', '入选', 'selected', 'chosen',
        '获奖', 'award', '提名', 'nominated'
    ]

    for line in lines:
        line_stripped = line.strip()
        if not line_stripped or len(line_stripped) < 4:
            continue
        line_lower = line_stripped.lower()
        for kw in achievement_keywords:
            if kw in line_lower:
                if line_stripped not in conclusions:
                    conclusions.append(line_stripped)
                break

    conclusion_indicators = [
        '总结', '结论', '建议', '效果', '成果', '收获',
        'summary', 'conclusion', 'result', 'outcome', 'takeaway',
        'successful', '成功', '显著', '复利', 'compound',
        'congratulations', '恭喜'
    ]

    for line in lines:
        line_stripped = line.strip()
        if not line_stripped or len(line_stripped) < 4:
            continue
        line_lower = line_stripped.lower()
        for kw in conclusion_indicators:
            if kw in line_lower and line_stripped not in conclusions:
                conclusions.append(line_stripped)
                break

    refined_items = []
    for dp in key_data_points:
        item = _refine_data_point(dp)
        if item and item not in refined_items:
            refined_items.append(item)
    for c in conclusions:
        item = _refine_conclusion(c)
        if item and item not in refined_items:
            refined_items.append(item)

    if not refined_items:
        for line in lines:
            line_stripped = line.strip()
            if not line_stripped or len(line_stripped) < 8:
                continue
            if re.search(r'\d', line_stripped):
                translated = _refine_data_point(line_stripped)
                if translated:
                    refined_items.append(translated)
                else:
                    refined_items.append(_translate_to_chinese(line_stripped) if not re.search(r'[\u4e00-\u9fff]', line_stripped) else line_stripped)
            elif any(kw in line_stripped.lower() for kw in ['重要', 'key', 'note', '注意', '关键', 'important', 'summary', 'conclusion']):
                translated = _translate_to_chinese(line_stripped) if not re.search(r'[\u4e00-\u9fff]', line_stripped) else line_stripped
                refined_items.append(translated)
            if len(refined_items) >= 10:
                break

    final_items = []
    seen = set()
    for item in refined_items:
        normalized = re.sub(r'\s+', '', item.lower())
        if normalized not in seen:
            seen.add(normalized)
            final_items.append(item)
        if len(final_items) >= 10:
            break

    return {
        "parsed_items": final_items,
        "raw_preview": full_text[:1000] + ("..." if len(full_text) > 1000 else ""),
        "total_chars": len(full_text)
    }


def _refine_data_point(text: str) -> str:
    """将数据点简化为简洁中文表达"""
    text = text.strip()
    if not text:
        return ""
    # 过滤掉太短或无意义的碎片
    if len(text) < 5:
        return ""
    # 过滤掉纯标点或纯空白
    if not re.search(r'[\w\u4e00-\u9fff]', text):
        return ""
    numbers = re.findall(r'[\d,]+(?:\.\d+)?', text)
    text_lower = text.lower()
    if re.search(r'[\u4e00-\u9fff]', text) and len(text) <= 80:
        return text
    if 'wishlist' in text_lower or '愿望单' in text_lower:
        for num in numbers:
            clean_num = num.replace(',', '')
            if clean_num.isdigit() and int(clean_num) > 100:
                date_match = re.search(r'(\d{1,2})[./](\d{1,2})', text)
                date_str = ""
                if date_match:
                    date_str = f"{date_match.group(1)}月{date_match.group(2)}日"
                result = f"{date_str}愿望单增长{num}" if date_str else f"愿望单增长{num}"
                if 'doubl' in text_lower or '翻倍' in text_lower or '2倍' in text_lower or '2x' in text_lower:
                    result += "，为平时的2倍"
                elif 'triple' in text_lower or '3倍' in text_lower or '3x' in text_lower:
                    result += "，为平时的3倍"
                elif re.search(r'(\d+)\s*x\b', text_lower):
                    multiplier = re.search(r'(\d+)\s*x\b', text_lower).group(1)
                    result += f"，为平时的{multiplier}倍"
                return result
    if any(kw in text_lower for kw in ['play', '游玩', 'player', 'concurrent', '同时在线']):
        for num in numbers:
            clean_num = num.replace(',', '')
            if clean_num.isdigit() and int(clean_num) > 100:
                if 'concurrent' in text_lower or '同时在线' in text_lower or 'peak' in text_lower:
                    return f"同时在线峰值{num}人"
                return f"游玩人数{num}"
    if any(kw in text_lower for kw in ['download', '下载', 'demo download']):
        for num in numbers:
            clean_num = num.replace(',', '')
            if clean_num.isdigit() and int(clean_num) > 50:
                return f"下载量达{num}"
    if any(kw in text_lower for kw in ['sales', 'sold', 'copies', '销量', '套']):
        for num in numbers:
            clean_num = num.replace(',', '')
            if clean_num.isdigit() and int(clean_num) > 50:
                return f"销量达{num}套"
    pct_match = re.search(r'([\d.]+)%', text)
    if pct_match:
        pct = pct_match.group(1)
        if any(kw in text_lower for kw in ['conversion', '转化']):
            return f"转化率{pct}%"
        elif any(kw in text_lower for kw in ['increase', 'growth', '增长', 'grew', 'rise']):
            return f"增长率{pct}%"
        elif any(kw in text_lower for kw in ['retention', '留存']):
            return f"留存率{pct}%"
    if any(kw in text_lower for kw in ['doubl', '翻倍']):
        date_match = re.search(r'(?:since|从|自)\s*(\d{1,2})[./](\d{1,2})', text, re.IGNORECASE)
        if date_match:
            return f"自{date_match.group(1)}月{date_match.group(2)}日起日增长翻倍"
        return "增长趋势翻倍"
    if not re.search(r'[\u4e00-\u9fff]', text):
        return _translate_to_chinese(text)
    if len(text) <= 80:
        return text
    return text[:80] + "..."


def _translate_to_chinese(text: str) -> str:
    """将英文数据文本翻译为简洁中文表达"""
    text = text.strip()
    if not text or len(text) < 3:
        return text
    text_lower = text.lower()
    numbers = re.findall(r'[\d,]+(?:\.\d+)?', text)
    pct_match = re.search(r'([\d.]+)%', text)

    # 先尝试基于关键词的语义翻译，生成完整中文句子
    num_str = numbers[0] if numbers else ""

    if 'wishlist' in text_lower or 'wishlists' in text_lower:
        if num_str:
            if 'daily' in text_lower or '日' in text_lower:
                return f"日均愿望单增长{num_str}"
            if 'total' in text_lower or '总' in text_lower:
                return f"愿望单总量达{num_str}"
            return f"愿望单数据：{num_str}"
        return "愿望单相关数据"

    if 'download' in text_lower or 'demo' in text_lower:
        if num_str:
            return f"下载量达{num_str}"
        if 'demo' in text_lower:
            return "试玩版下载数据"

    if any(kw in text_lower for kw in ['player', 'players', 'concurrent', 'played']):
        if num_str:
            if 'concurrent' in text_lower or 'peak' in text_lower:
                return f"同时在线峰值{num_str}人"
            return f"玩家数达{num_str}"

    if any(kw in text_lower for kw in ['sale', 'sold', 'copies', 'revenue']):
        if num_str:
            if 'revenue' in text_lower:
                return f"收入达{num_str}"
            return f"销量达{num_str}"

    if pct_match:
        pct = pct_match.group(1)
        if any(kw in text_lower for kw in ['conversion', 'convert']):
            return f"转化率{pct}%"
        elif any(kw in text_lower for kw in ['increase', 'growth', 'grew', 'rise', 'grow']):
            return f"增长{pct}%"
        elif any(kw in text_lower for kw in ['retention', 'retain']):
            return f"留存率{pct}%"
        return f"关键指标{pct}%"

    if any(kw in text_lower for kw in ['top', 'ranking', 'chart', 'rank']):
        rank_match = re.search(r'top[- ]?(\d+)', text_lower)
        if rank_match:
            return f"进入Top-{rank_match.group(1)}排名"
        return "获得榜单排名成绩"

    if any(kw in text_lower for kw in ['doubled', 'double', 'tripled', 'triple']):
        if 'triple' in text_lower:
            return "增长3倍"
        return "增长翻倍"

    if any(kw in text_lower for kw in ['next fest', 'spring sale', 'summer sale', 'autumn sale', 'winter sale']):
        event_map = {
            'next fest': '新品节', 'spring sale': '春季特卖',
            'summer sale': '夏季特卖', 'autumn sale': '秋季特卖',
            'winter sale': '冬季特卖'
        }
        for eng, chn in event_map.items():
            if eng in text_lower:
                if num_str:
                    return f"{chn}期间数据：{num_str}"
                return f"参与{chn}活动"

    # 如果无法语义翻译，进行词汇替换
    translations = {
        'wishlist': '愿望单', 'wishlists': '愿望单',
        'download': '下载', 'downloads': '下载',
        'player': '玩家', 'players': '玩家',
        'sales': '销量', 'sold': '售出', 'copies': '套',
        'daily': '日均', 'weekly': '周均', 'monthly': '月均',
        'increase': '增长', 'growth': '增长', 'grew': '增长了',
        'decrease': '下降', 'drop': '下降',
        'peak': '峰值', 'total': '总计', 'average': '平均',
        'revenue': '收入', 'conversion': '转化率',
        'retention': '留存率', 'engagement': '参与度',
        'impressions': '曝光次数', 'clicks': '点击次数',
        'views': '浏览量', 'page views': '页面浏览量',
        'unique visitors': '独立访客',
        'demo': '试玩版', 'next fest': '新品节',
        'spring sale': '春季特卖', 'summer sale': '夏季特卖',
        'autumn sale': '秋季特卖', 'winter sale': '冬季特卖',
        'day one': '首日', 'first day': '首日', 'first week': '首周',
        'launch': '发布', 'release': '发布',
        'region': '区域', 'global': '全球', 'worldwide': '全球',
        'china': '中国', 'asia': '亚洲', 'europe': '欧洲',
        'north america': '北美', 'japan': '日本', 'korea': '韩国',
        'most played': '游玩人数最多', 'most wished': '最受期待',
        'top seller': '最畅销', 'trending': '热门',
        'featured': '推荐展示', 'highlighted': '重点推荐',
        'doubled': '翻倍', 'tripled': '增长3倍',
        'achieve': '达成', 'achieved': '达成',
        'successful': '成功', 'success': '成功',
    }
    result = text
    sorted_keys = sorted(translations.keys(), key=len, reverse=True)
    for eng, chn in [(k, translations[k]) for k in sorted_keys]:
        pattern = re.compile(re.escape(eng), re.IGNORECASE)
        result = pattern.sub(chn, result)
    result = re.sub(r'\s+', ' ', result).strip()
    # 清理残留的英文连接词和介词
    result = re.sub(r'\b(and|or|the|a|an|of|in|on|at|to|for|with|from|by|is|was|are|were|has|have|had|got|get)\b', '', result, flags=re.IGNORECASE)
    result = re.sub(r'\s+', ' ', result).strip()
    result = re.sub(r'^\s*[.,;:]\s*', '', result)
    result = re.sub(r'\s*[.,;:]\s*$', '', result)

    chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', result))
    total_chars = len(result.replace(' ', ''))
    if total_chars > 0 and chinese_chars / total_chars < 0.15 and numbers:
        # 翻译效果不好，用简洁的中文概括
        if 'wishlist' in text_lower:
            return f"愿望单相关数据：{num_str}" if num_str else "愿望单相关数据"
        elif 'download' in text_lower:
            return f"下载量：{num_str}" if num_str else "下载相关数据"
        elif 'player' in text_lower or 'play' in text_lower:
            return f"玩家数：{num_str}" if num_str else "玩家相关数据"
        elif 'sale' in text_lower or 'sold' in text_lower:
            return f"销售数据：{num_str}" if num_str else "销售相关数据"
        else:
            return f"关键数据：{num_str}" if num_str else ""
    if len(result) > 80:
        return result[:80] + "..."
    return result


def _refine_conclusion(text: str) -> str:
    """将结论简化为简洁中文表达"""
    text = text.strip()
    if not text:
        return ""
    text_lower = text.lower()
    if 'most played' in text_lower or '游玩人数最多' in text_lower:
        if 'demo' in text_lower:
            region_parts = []
            if 'china' in text_lower or '中国' in text_lower:
                region_parts.append("中国")
            if 'overseas' in text_lower or '海外' in text_lower or 'global' in text_lower:
                region_parts.append("海外")
            region = f"（{'+'.join(region_parts)}）" if region_parts else ""
            return f"进入游玩人数最多Demo榜单{region}"
        return "进入游玩人数最多榜单"
    if 'most wished' in text_lower or '最受期待' in text_lower:
        return "进入最受期待榜单"
    if 'top seller' in text_lower or '最畅销' in text_lower:
        return "进入最畅销榜单"
    if 'trending' in text_lower or '热门' in text_lower:
        return "进入热门趋势榜"
    if 'new release' in text_lower and ('top' in text_lower or 'chart' in text_lower):
        return "进入新品畅销榜"
    if 'highlighted' in text_lower or 'featured' in text_lower:
        return "获平台官方推荐展示"
    if 'front page' in text_lower or '首页' in text_lower:
        return "获得商店首页推荐"
    if 'successful' in text_lower or '成功' in text_lower:
        if 'next fest' in text_lower or '新品节' in text_lower:
            return "新品节参与效果显著"
        if 'launch' in text_lower or '发布' in text_lower:
            return "游戏发布取得成功"
        return "活动参与效果显著"
    if 'compound' in text_lower or '复利' in text_lower:
        return "营销推广产生复利效应"
    if 'milestone' in text_lower or '里程碑' in text_lower:
        return "达成重要里程碑"
    if 'award' in text_lower or '获奖' in text_lower or 'nominated' in text_lower:
        return "获得游戏奖项/提名"
    if re.search(r'[\u4e00-\u9fff]', text):
        return text[:80]
    translated = _translate_to_chinese(text)
    chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', translated))
    if chinese_chars > 0:
        return translated
    if 'chart' in text_lower or 'ranking' in text_lower or 'top' in text_lower:
        return "获得榜单排名成绩"
    if 'event' in text_lower or 'festival' in text_lower:
        return "参与平台活动获得成效"
    if len(text) <= 60:
        return text
    return text[:60] + "..."


class ReportParseRequest(BaseModel):
    content: str
    filename: str = ""


@app.post("/api/documents/parse-report")
async def parse_report(data: ReportParseRequest):
    """解析报告内容，提炼关键数据和结论"""
    if not data.content or not data.content.strip():
        return {"code": -1, "msg": "报告内容为空"}
    result = parse_report_content(data.content, data.filename)
    _cleanup_expired_cache()
    temp_id = str(uuid.uuid4())
    _report_temp_cache[temp_id] = {
        "original_content": data.content,
        "file_type": "txt",
        "pdf_data": None,
        "created_at": datetime.now()
    }
    result["temp_id"] = temp_id
    return {"code": 0, "data": result, "msg": "解析完成"}


@app.post("/api/documents/upload-report")
async def upload_report_document(
    platform_id: Optional[int] = Form(None),
    game_name: Optional[str] = Form(None),
    parsed_content: str = Form(...),
    filename: str = Form("内部项目报告"),
    temp_id: str = Form("")
):
    """上传经用户确认的报告"""
    try:
        original_content = ""
        file_type = "txt"
        pdf_data_bytes = None
        if temp_id and temp_id in _report_temp_cache:
            cached = _report_temp_cache[temp_id]
            original_content = cached.get("original_content", "")
            file_type = cached.get("file_type", "txt")
            pdf_data_bytes = cached.get("pdf_data")
            del _report_temp_cache[temp_id]

        result = wxcloud_db.add_document(
            platform_id=platform_id,
            filename=filename,
            content=parsed_content,
            doc_type="report",
            game_name=game_name,
            file_type=file_type,
            original_content=original_content,
            pdf_data=pdf_data_bytes,
            status="pending"
        )
        doc_id = result.get("doc_id", 0)

        wxcloud_db.add_update_log(
            platform_id, f"上传内部项目报告: {filename}",
            f"游戏: {game_name or '未指定'}, 文档ID: {doc_id}")

        return {"code": 0, "data": {"id": doc_id, "filename": filename}, "msg": "报告上传成功"}
    except Exception as e:
        return {"code": -1, "msg": f"上传失败: {str(e)}"}


def _parse_file_content(content: bytes, filename: str):
    """解析上传文件内容，返回 (text_content, file_type)"""
    filename_lower = (filename or "").lower()
    file_type = "txt"

    if filename_lower.endswith(".docx"):
        file_type = "docx"
        try:
            import io
            from docx import Document
            from docx.table import Table as DocxTable
            from docx.text.paragraph import Paragraph as DocxParagraph
            from docx.oxml.ns import qn
            doc = Document(io.BytesIO(content))
            elements = []
            for child in doc.element.body:
                if child.tag == qn('w:p'):
                    para = DocxParagraph(child, doc.element.body)
                    text = para.text.strip()
                    if text:
                        elements.append(text)
                elif child.tag == qn('w:tbl'):
                    table = DocxTable(child, doc.element.body)
                    elements.append("<!--TABLE_START-->")
                    for row in table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells]
                        elements.append(" | ".join(row_texts))
                    elements.append("<!--TABLE_END-->")
            text_content = "\n".join(elements)
        except Exception as e:
            return f"(无法解析docx文件: {str(e)})", file_type
    elif filename_lower.endswith(".doc"):
        return "(不支持旧版.doc格式，请转换为.docx后重新上传)", "doc"
    elif filename_lower.endswith(".pptx"):
        file_type = "pptx"
        try:
            import io
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation(io.BytesIO(content))
            elements = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_texts = []
                slide_tables = []
                # 按位置排序shapes（从上到下，从左到右）
                sorted_shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))
                for shape in sorted_shapes:
                    if shape.has_text_frame:
                        # 将同一个shape的所有段落合并
                        shape_lines = []
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text and len(text) >= 2:
                                shape_lines.append(text)
                        if shape_lines:
                            # 如果shape内多行文本，合并为一个文本块
                            combined = "\n".join(shape_lines)
                            slide_texts.append(combined)
                    if shape.has_table:
                        table = shape.table
                        table_lines = ["<!--TABLE_START-->"]
                        for row in table.rows:
                            row_texts = [cell.text.strip() for cell in row.cells]
                            table_lines.append(" | ".join(row_texts))
                        table_lines.append("<!--TABLE_END-->")
                        slide_tables.append("\n".join(table_lines))

                # 只有有实质内容的幻灯片才输出
                if slide_texts or slide_tables:
                    elements.append(f"--- 第 {slide_idx} 页 ---")
                    # 合并过短的相邻文本块（可能是被拆分的句子）
                    merged_texts = []
                    buffer = ""
                    for t in slide_texts:
                        # 如果当前文本很短且不含换行，尝试与buffer合并
                        if len(t) < 15 and "\n" not in t and buffer:
                            buffer = buffer + " " + t
                        else:
                            if buffer:
                                merged_texts.append(buffer)
                            buffer = t
                    if buffer:
                        merged_texts.append(buffer)

                    for t in merged_texts:
                        elements.append(t)
                    for t in slide_tables:
                        elements.append(t)

                    # 提取备注
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text and len(notes_text) >= 3:
                            elements.append(f"[备注] {notes_text}")
                    elements.append("")
            text_content = "\n".join(elements)
        except Exception as e:
            return f"(无法解析PPTX文件: {str(e)})", file_type
    elif filename_lower.endswith(".ppt"):
        file_type = "ppt"
        return "(不支持旧版.ppt格式，请转换为.pptx后重新上传)", file_type
    elif filename_lower.endswith(".xlsx") or filename_lower.endswith(".xls"):
        file_type = "xlsx" if filename_lower.endswith(".xlsx") else "xls"
        try:
            import io
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            elements = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                elements.append(f"--- 工作表: {sheet_name} ---")
                elements.append("<!--TABLE_START-->")
                for row in ws.iter_rows(values_only=True):
                    row_texts = [(str(cell) if cell is not None else "") for cell in row]
                    if any(t.strip() for t in row_texts):
                        elements.append(" | ".join(row_texts))
                elements.append("<!--TABLE_END-->")
                elements.append("")
            wb.close()
            text_content = "\n".join(elements)
        except Exception as e:
            return f"(无法解析Excel文件: {str(e)})", file_type
    elif filename_lower.endswith(".pdf"):
        file_type = "pdf"
        try:
            import io, pdfplumber
            pdf = pdfplumber.open(io.BytesIO(content))
            elements = []
            for page in pdf.pages:
                tables = page.extract_tables()
                if tables:
                    page_text = page.extract_text(layout=False)
                    table_texts_set = set()
                    for table in tables:
                        for row in table:
                            for cell in row:
                                if cell:
                                    table_texts_set.add(cell.strip())
                    if page_text:
                        for line in page_text.split('\n'):
                            ls = line.strip()
                            if ls and ls not in table_texts_set:
                                words = ls.split()
                                in_table_count = sum(1 for w in words if w.strip() in table_texts_set)
                                if len(words) > 0 and in_table_count / len(words) < 0.5:
                                    elements.append(ls)
                    for table in tables:
                        elements.append("<!--TABLE_START-->")
                        for row in table:
                            row_texts = [(cell or "").strip().replace('\n', ' ') for cell in row]
                            elements.append(" | ".join(row_texts))
                        elements.append("<!--TABLE_END-->")
                else:
                    page_text = page.extract_text()
                    if page_text:
                        for line in page_text.split('\n'):
                            ls = line.strip()
                            if ls:
                                elements.append(ls)
                elements.append("")
            pdf.close()
            text_content = "\n".join(elements)
        except Exception as e:
            return f"(无法解析PDF文件: {str(e)})", file_type
    elif filename_lower.endswith(".md") or filename_lower.endswith(".markdown"):
        file_type = "md"
        try:
            text_content = content.decode("utf-8")
        except UnicodeDecodeError:
            text_content = content.decode("gbk", errors="ignore")
    elif filename_lower.endswith(".html") or filename_lower.endswith(".htm"):
        file_type = "html"
        try:
            import re as _re
            raw_html = content.decode("utf-8", errors="ignore")
            # 移除script和style标签及内容
            clean = _re.sub(r'<(script|style)[^>]*>.*?</\1>', '', raw_html, flags=_re.DOTALL | _re.IGNORECASE)
            # 移除HTML标签
            clean = _re.sub(r'<[^>]+>', '\n', clean)
            # 清理多余空行
            lines = [line.strip() for line in clean.split('\n') if line.strip()]
            text_content = "\n".join(lines)
        except Exception as e:
            return f"(无法解析HTML文件: {str(e)})", file_type
    elif filename_lower.endswith(".rtf"):
        file_type = "rtf"
        try:
            import re as _re
            raw = content.decode("utf-8", errors="ignore")
            # 简单RTF文本提取：移除RTF控制字和花括号
            clean = _re.sub(r'\\[a-z]+\d*\s?', '', raw)
            clean = _re.sub(r'[{}]', '', clean)
            clean = _re.sub(r'\\\*.*?\n', '', clean)
            lines = [line.strip() for line in clean.split('\n') if line.strip()]
            text_content = "\n".join(lines)
        except Exception as e:
            return f"(无法解析RTF文件: {str(e)})", file_type
    elif filename_lower.endswith(".csv"):
        file_type = "csv"
        try:
            import io, csv
            raw_text = content.decode("utf-8", errors="ignore")
            reader = csv.reader(io.StringIO(raw_text))
            elements = ["<!--TABLE_START-->"]
            for row in reader:
                elements.append(" | ".join(row))
            elements.append("<!--TABLE_END-->")
            text_content = "\n".join(elements)
        except Exception as e:
            return f"(无法解析CSV文件: {str(e)})", file_type
    else:
        try:
            text_content = content.decode("utf-8")
        except UnicodeDecodeError:
            try:
                text_content = content.decode("gbk", errors="ignore")
            except Exception:
                text_content = content.decode("latin-1", errors="ignore")

    return text_content, file_type


@app.post("/api/documents/parse-report-file")
async def parse_report_file(file: UploadFile = File(...)):
    """上传文件并解析报告内容"""
    content = await file.read()
    text_content, file_type = _parse_file_content(content, file.filename)
    if isinstance(text_content, tuple):
        text_content = text_content[0]
    if not text_content or not text_content.strip():
        return {"code": -1, "msg": "文件内容为空"}
    result = parse_report_content(text_content, file.filename or "")
    result["filename"] = file.filename or "未命名文件"
    result["file_type"] = file_type
    _cleanup_expired_cache()
    temp_id = str(uuid.uuid4())
    _report_temp_cache[temp_id] = {
        "original_content": text_content,
        "file_type": file_type,
        "pdf_data": content if file_type in ("pdf", "docx", "doc", "pptx", "ppt", "xlsx", "xls", "csv", "html", "htm", "rtf", "md") else None,
        "created_at": datetime.now()
    }
    result["temp_id"] = temp_id
    result["original_preview"] = text_content[:1500] + ("..." if len(text_content) > 1500 else "")
    return {"code": 0, "data": result, "msg": "解析完成"}


@app.post("/api/documents/upload")
async def upload_document(
    file: UploadFile = File(...),
    platform_id: Optional[int] = Form(None),
    doc_type: str = Form("policy"),
    game_name: Optional[str] = Form(None)
):
    """上传文档"""
    content = await file.read()
    text_content, file_type = _parse_file_content(content, file.filename)
    if isinstance(text_content, tuple):
        text_content = text_content[0]
    # 保存原始文件二进制数据（PDF/DOCX/PPTX/XLSX等）
    pdf_data = content if file_type in ("pdf", "docx", "doc", "pptx", "ppt", "xlsx", "xls", "csv", "html", "htm", "rtf", "md") else None
    try:
        result = wxcloud_db.add_document(
            platform_id=platform_id,
            filename=file.filename,
            content=text_content,
            doc_type=doc_type,
            game_name=game_name,
            file_type=file_type,
            original_content=text_content,
            pdf_data=pdf_data,
            status="pending"
        )
        doc_id = result.get("doc_id", 0)
        wxcloud_db.add_update_log(
            platform_id, f"上传文档: {file.filename}",
            f"文档类型: {doc_type}, 文档ID: {doc_id}")
        return {"code": 0, "data": {"id": doc_id, "filename": file.filename}, "msg": "上传成功"}
    except Exception as e:
        return {"code": -1, "msg": f"上传失败: {str(e)}"}


@app.post("/api/documents/upload-text")
async def upload_text_document(doc: DocUpload):
    """上传文本内容"""
    try:
        result = wxcloud_db.add_document(
            platform_id=doc.platform_id,
            filename=doc.filename,
            content=doc.content,
            doc_type=doc.doc_type,
            game_name=doc.game_name,
            status="pending"
        )
        doc_id = result.get("doc_id", 0)
        wxcloud_db.add_update_log(
            doc.platform_id, f"上传文本: {doc.filename}",
            f"文档类型: {doc.doc_type}, 文档ID: {doc_id}")
        return {"code": 0, "data": {"id": doc_id, "filename": doc.filename}, "msg": "上传成功"}
    except Exception as e:
        return {"code": -1, "msg": f"上传失败: {str(e)}"}


@app.get("/api/documents")
async def get_documents(status: Optional[str] = None):
    """获取上传的文档列表"""
    try:
        docs = wxcloud_db.get_documents(status=status)
        return {"code": 0, "data": docs}
    except Exception as e:
        return {"code": -1, "msg": f"加载文档失败: {str(e)}"}


@app.get("/api/documents/{doc_id}")
async def get_document_detail(doc_id: int):
    """获取单个文档详情"""
    try:
        doc = wxcloud_db.get_document_by_id(doc_id)
        if not doc:
            return {"code": -1, "msg": "文档不存在"}
        # 移除大字段减少传输
        doc.pop("pdf_data", None)
        return {"code": 0, "data": doc}
    except Exception as e:
        return {"code": -1, "msg": f"加载文档失败: {str(e)}"}


@app.get("/api/documents/{doc_id}/pdf")
async def get_document_pdf(doc_id: int):
    """获取文档的原始PDF文件"""
    try:
        with wxcloud_db.get_db() as db:
            cursor = db.cursor()
            cursor.execute(
                "SELECT filename, pdf_data, file_type FROM uploaded_docs WHERE id=%s",
                (doc_id,)
            )
            doc = cursor.fetchone()
            cursor.close()
        
        if not doc:
            return JSONResponse(status_code=404, content={"code": -1, "msg": "文档不存在"})
        
        if doc.get("file_type") != "pdf" or not doc.get("pdf_data"):
            return JSONResponse(status_code=404, content={"code": -1, "msg": "该文档没有PDF原件"})
        
        filename = doc.get("filename", "document.pdf")
        if not filename.lower().endswith(".pdf"):
            filename += ".pdf"
        
        return Response(
            content=doc["pdf_data"],
            media_type="application/pdf",
            headers={
                "Content-Disposition": f"inline; filename=\"{filename}\"",
                "Cache-Control": "max-age=3600"
            }
        )
    except Exception as e:
        return JSONResponse(status_code=500, content={"code": -1, "msg": f"获取PDF失败: {str(e)}"})


@app.get("/api/documents/{doc_id}/original-file")
async def get_document_original_file(doc_id: int):
    """获取文档的原始上传文件（PDF/DOCX等）"""
    try:
        with wxcloud_db.get_db() as db:
            cursor = db.cursor()
            cursor.execute(
                "SELECT filename, pdf_data, file_type FROM uploaded_docs WHERE id=%s",
                (doc_id,)
            )
            doc = cursor.fetchone()
            cursor.close()

        if not doc:
            return JSONResponse(status_code=404, content={"code": -1, "msg": "文档不存在"})

        file_data = doc.get("pdf_data")
        file_type = doc.get("file_type", "txt")
        filename = doc.get("filename", "document")

        if not file_data:
            return JSONResponse(status_code=404, content={"code": -1, "msg": "该文档没有原始文件"})

        # 根据文件类型设置MIME类型
        mime_map = {
            "pdf": "application/pdf",
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "doc": "application/msword",
            "txt": "text/plain",
            "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "csv": "text/csv",
        }
        mime_type = mime_map.get(file_type, "application/octet-stream")

        # 确保文件名有正确的扩展名
        ext_map = {"pdf": ".pdf", "docx": ".docx", "doc": ".doc", "txt": ".txt", "xlsx": ".xlsx", "csv": ".csv"}
        expected_ext = ext_map.get(file_type, "")
        if expected_ext and not filename.lower().endswith(expected_ext):
            filename += expected_ext

        return Response(
            content=file_data,
            media_type=mime_type,
            headers={
                "Content-Disposition": f"inline; filename=\"{filename}\"",
                "Cache-Control": "max-age=3600"
            }
        )
    except Exception as e:
        return JSONResponse(status_code=500, content={"code": -1, "msg": f"获取原始文件失败: {str(e)}"})


@app.get("/api/documents/{doc_id}/extract-keypoints")
async def extract_document_keypoints(doc_id: int):
    """提炼文档关键点"""
    try:
        doc = wxcloud_db.get_document_by_id(doc_id)
        if not doc:
            return {"code": -1, "msg": "文档不存在"}

        content = doc.get("content", "") or ""
        original_content = doc.get("original_content", "") or ""
        # 优先使用原始内容进行分析
        text_to_analyze = original_content if original_content else content
        if not text_to_analyze.strip():
            return {"code": -1, "msg": "文档内容为空，无法提炼"}

        filename = doc.get("filename", "")
        doc_type = doc.get("doc_type", "other")

        keypoints = _extract_keypoints_from_text(text_to_analyze, filename, doc_type)

        return {
            "code": 0,
            "data": {
                "doc_id": doc_id,
                "filename": filename,
                "keypoints": keypoints,
                "total_chars": len(text_to_analyze)
            },
            "msg": "关键点提炼完成"
        }
    except Exception as e:
        return {"code": -1, "msg": f"提炼失败: {str(e)}"}


def _extract_keypoints_from_text(text: str, filename: str = "", doc_type: str = "other") -> list:
    """从文本中提炼关键点，返回结构化的关键点列表"""
    lines = text.strip().split('\n')
    keypoints = []
    seen = set()

    # 1. 提取数字相关的关键数据
    number_patterns = [
        (r'[^\n]*?(?:愿望单|wishlist)[^\n]*?[\d,]+[^\n]*', '📊 数据指标'),
        (r'[^\n]*?(?:下载量|downloads?|销量|sales)[^\n]*?[\d,]+[^\n]*', '📊 数据指标'),
        (r'[^\n]*?(?:玩家|players?|用户|users?)[^\n]*?[\d,]+[^\n]*', '📊 数据指标'),
        (r'[^\n]*?[\d,.]+%[^\n]*?(?:转化|conversion|增长|growth|留存|retention)[^\n]*', '📈 增长趋势'),
        (r'[^\n]*?(?:增长|增加|提升|翻倍|doubl|triple|grew|increase|boost)[^\n]*?[\d,.]+[^\n]*', '📈 增长趋势'),
        (r'[^\n]*?(?:收入|revenue|分成|佣金)[^\n]*?[\d,.]+[^\n]*', '💰 收入数据'),
    ]

    for pattern, category in number_patterns:
        matches = re.finditer(pattern, text, re.IGNORECASE)
        for m in matches:
            raw = m.group(0).strip()
            if len(raw) < 5 or len(raw) > 200:
                continue
            normalized = re.sub(r'\s+', '', raw.lower())
            if normalized not in seen:
                seen.add(normalized)
                refined = _refine_data_point(raw)
                if refined:
                    keypoints.append({"text": refined, "category": category, "importance": "high"})

    # 2. 提取时间/日期相关信息
    date_patterns = [
        r'[^\n]*?(?:\d{4}年\d{1,2}月|\d{1,2}月\d{1,2}日)[^\n]*?(?:开始|结束|截止|上线|发布|活动)[^\n]*',
        r'[^\n]*?(?:开始|结束|截止|上线|发布|活动)[^\n]*?(?:\d{4}年\d{1,2}月|\d{1,2}月\d{1,2}日)[^\n]*',
        r'[^\n]*?(?:deadline|launch|release|start|end)[^\n]*?\d{4}[/-]\d{1,2}[/-]\d{1,2}[^\n]*',
    ]
    for pattern in date_patterns:
        matches = re.finditer(pattern, text, re.IGNORECASE)
        for m in matches:
            raw = m.group(0).strip()
            if len(raw) < 8 or len(raw) > 150:
                continue
            normalized = re.sub(r'\s+', '', raw.lower())
            if normalized not in seen:
                seen.add(normalized)
                keypoints.append({"text": raw[:100], "category": "📅 时间节点", "importance": "medium"})

    # 3. 提取政策/规则要点
    policy_keywords = [
        '要求', '必须', '不得', '禁止', '规定', '条件', '标准', '限制',
        'require', 'must', 'shall', 'prohibit', 'mandatory', 'criteria',
        '分成', '比例', '费用', '价格', '折扣',
    ]
    for line in lines:
        line_stripped = line.strip()
        if not line_stripped or len(line_stripped) < 6:
            continue
        line_lower = line_stripped.lower()
        for kw in policy_keywords:
            if kw in line_lower:
                normalized = re.sub(r'\s+', '', line_lower)
                if normalized not in seen:
                    seen.add(normalized)
                    keypoints.append({"text": line_stripped[:120], "category": "📋 政策规则", "importance": "high"})
                break

    # 4. 提取结论/成果
    conclusion_keywords = [
        '总结', '结论', '建议', '效果', '成果', '收获', '成功',
        'summary', 'conclusion', 'result', 'outcome', 'recommendation',
        '榜单', '排名', 'top', 'ranking', '推荐', '入选',
    ]
    for line in lines:
        line_stripped = line.strip()
        if not line_stripped or len(line_stripped) < 6:
            continue
        line_lower = line_stripped.lower()
        for kw in conclusion_keywords:
            if kw in line_lower:
                normalized = re.sub(r'\s+', '', line_lower)
                if normalized not in seen:
                    seen.add(normalized)
                    refined = _refine_conclusion(line_stripped)
                    if refined:
                        keypoints.append({"text": refined, "category": "🏆 成果结论", "importance": "medium"})
                break

    # 5. 如果关键点太少，提取标题/重点行
    if len(keypoints) < 3:
        for line in lines:
            line_stripped = line.strip()
            if not line_stripped or len(line_stripped) < 4:
                continue
            # 识别标题行（短行、以数字/符号开头等）
            is_heading = (len(line_stripped) < 50 and not line_stripped.endswith('。') and not line_stripped.endswith('.'))
            has_number = bool(re.search(r'\d', line_stripped))
            if is_heading or has_number:
                normalized = re.sub(r'\s+', '', line_stripped.lower())
                if normalized not in seen:
                    seen.add(normalized)
                    keypoints.append({"text": line_stripped[:100], "category": "📝 要点摘要", "importance": "low"})
            if len(keypoints) >= 15:
                break

    # 按重要性排序
    importance_order = {"high": 0, "medium": 1, "low": 2}
    keypoints.sort(key=lambda x: importance_order.get(x.get("importance", "low"), 2))

    return keypoints[:20]


@app.delete("/api/documents/{doc_id}")
async def delete_document_api(doc_id: int):
    """删除文档"""
    try:
        wxcloud_db.delete_document(doc_id)
        return {"code": 0, "msg": "删除成功"}
    except Exception as e:
        return {"code": -1, "msg": f"删除失败: {str(e)}"}


@app.post("/api/strategies/update")
async def update_strategy(data: StrategyUpdate):
    """更新或新增策略"""
    try:
        # 查找现有策略
        where = f'{{platform_id: {data.platform_id}, category: "{data.category}"}}'
        result = wxcloud_db.db_query("strategies", where=where, limit=1)
        existing_list = result.get("data", [])

        if existing_list:
            existing = existing_list[0]
            new_version = existing.get("version", 1) + 1
            wxcloud_db.db_update_by_id("strategies", existing["id"], {
                "title": data.title,
                "content": data.content,
                "version": new_version,
                "updated_at": wxcloud_db._now_str()
            })
        else:
            new_strategy = {
                "platform_id": data.platform_id,
                "category": data.category,
                "title": data.title,
                "content": data.content,
                "version": 1,
                "created_at": wxcloud_db._now_str(),
                "updated_at": wxcloud_db._now_str()
            }
            wxcloud_db.db_add("strategies", new_strategy)

        # 标记相关文档为已应用
        pending_where = f'{{platform_id: {data.platform_id}, status: "pending"}}'
        wxcloud_db.db_update("uploaded_docs", pending_where, {"status": "applied", "updated_at": wxcloud_db._now_str()})

        wxcloud_db.add_update_log(data.platform_id, f"更新策略: {data.category}", f"标题: {data.title}")
        return {"code": 0, "msg": "策略更新成功"}
    except Exception as e:
        return {"code": -1, "msg": f"更新失败: {str(e)}"}


@app.post("/api/strategies/batch-update")
async def batch_update_strategies(platform_id: int = Form(...)):
    """基于已上传的文档批量更新策略"""
    try:
        # 获取所有pending文档
        all_pending_docs = wxcloud_db.get_pending_docs()

        if not all_pending_docs:
            return {"code": 1, "msg": "没有待处理的新文档信息"}

        # 分类文档
        normal_docs = [d for d in all_pending_docs if d.get("doc_type") != "report"
                      and (d.get("platform_id") == platform_id or d.get("platform_id") is None)]
        report_docs = [d for d in all_pending_docs if d.get("doc_type") == "report"]
        pending_docs = normal_docs + report_docs

        if not pending_docs:
            return {"code": 1, "msg": "没有待处理的新文档信息"}

        # 获取当前平台策略
        current_strategies = wxcloud_db.get_strategies_by_platform(platform_id)

        # 获取平台名称
        platforms = wxcloud_db.get_all_platforms()
        platform_name = ""
        for p in platforms:
            if p["id"] == platform_id:
                platform_name = p["name"].lower()
                break

        # 平台关键词
        platform_match_keywords = {
            "steam": ["steam", "valve", "steamworks", "steam deck", "deck", "wishlist", "next fest", "新品节", "playtest"],
            "epic": ["epic", "epic games", "eos", "mega sale", "虚幻引擎", "unreal", "first run", "now on epic"],
            "xbox": ["xbox", "game pass", "gamepass", "xgp", "microsoft", "微软", "showcase", "developer direct", "series s", "series x", "deals with gold", "xbox live"],
            "playstation": ["playstation", "ps5", "ps4", "ps plus", "psn", "sony", "索尼", "dualsense", "state of play", "trc", "activity card", "dealmania", "days of play"]
        }

        current_platform_kws = []
        for pk, kws in platform_match_keywords.items():
            if pk in platform_name.replace(" ", ""):
                current_platform_kws = kws
                break

        # 活动关键词映射
        activity_keywords_map = {
            "next fest": ["next fest", "新品节", "nextfest"],
            "playtest": ["playtest", "封闭测试", "play test"],
            "spring sale": ["spring sale", "春季特卖", "春促"],
            "summer sale": ["summer sale", "夏季特卖", "夏促"],
            "autumn sale": ["autumn sale", "秋季特卖", "秋促"],
            "winter sale": ["winter sale", "冬季特卖", "冬促"],
            "daily deal": ["daily deal", "每日特惠"],
            "mega sale": ["mega sale", "mega特卖"],
            "game pass": ["game pass", "gamepass", "xgp"],
            "ps plus": ["ps plus", "ps+"],
            "days of play": ["days of play"],
            "state of play": ["state of play"],
            "showcase": ["showcase", "展示会", "games showcase"],
            "developer direct": ["developer direct"],
            "free play days": ["free play days", "免费试玩日", "free play"],
            "gamescom": ["gamescom", "科隆游戏展"],
            "deck": ["steam deck", "deck verified"],
            "black friday": ["black friday", "黑五"],
            "dealmania": ["dealmania"],
            "spotlight sale": ["spotlight sale", "聚光灯特卖", "spotlight"],
            "first run": ["first run"],
        }

        strategy_contents = {}
        for strategy in current_strategies:
            sid = strategy["id"]
            strategy_contents[sid] = strategy.get("content", {})

        report_assignments = {}

        if report_docs:
            for doc in report_docs:
                doc_text = doc.get("content", "") or ""
                doc_filename = doc.get("filename", "") or ""
                doc_content_lower = (doc_text + " " + doc_filename).lower()
                game_name = doc.get("game_name") or ""

                platform_relevant = False
                if doc.get("platform_id") == platform_id:
                    platform_relevant = True
                else:
                    for kw in current_platform_kws:
                        if kw in doc_content_lower:
                            platform_relevant = True
                            break
                if not platform_relevant:
                    continue

                activity_scores = {}
                for activity_name, activity_kws in activity_keywords_map.items():
                    score = 0
                    for akw in activity_kws:
                        count = doc_content_lower.count(akw)
                        if count > 0:
                            score += count * len(akw)
                    if score > 0:
                        activity_scores[activity_name] = score

                global_best = None

                if activity_scores:
                    primary_activity = max(activity_scores, key=activity_scores.get)
                    primary_activity_kws = activity_keywords_map.get(primary_activity, [])

                    for strategy in current_strategies:
                        sid = strategy["id"]
                        content = strategy_contents.get(sid, {})
                        for section in content.get("sections", []):
                            sec_title_lower = section.get("title", "").lower()
                            sec_title_score = 0
                            for akw in primary_activity_kws:
                                if akw in sec_title_lower:
                                    sec_title_score += len(akw) * 10
                            if sec_title_score == 0:
                                continue
                            for item in section.get("items", []):
                                item_text = (item.get("name", "") + " " + item.get("desc", "")).lower()
                                item_score = sec_title_score
                                for akw in primary_activity_kws:
                                    if akw in item_text:
                                        item_score += len(akw) * 2
                                data_bonus = {"转化率": 50, "参考数据": 50, "数据统计": 45, "📊": 40, "统计": 35, "效果数据": 40, "数据": 20}
                                for bkw, bscore in data_bonus.items():
                                    if bkw in item_text:
                                        item_score += bscore
                                if re.search(r'\d{1,2}月\d{1,2}日', item.get("name", "")):
                                    item_score -= 30
                                if global_best is None or item_score > global_best["score"]:
                                    global_best = {"strategy_id": sid, "strategy": strategy, "section": section, "item": item, "score": item_score}

                if not global_best:
                    doc_keywords = set()
                    fn_words = re.findall(r'[\u4e00-\u9fff]+|[a-zA-Z]{2,}', doc_filename)
                    for w in fn_words:
                        doc_keywords.add(w.lower())
                    content_words = re.findall(r'[\u4e00-\u9fff]{2,}|[a-zA-Z]{3,}', doc_text[:300])
                    for w in content_words:
                        doc_keywords.add(w.lower())
                    stop_words = {"the", "and", "for", "with", "from", "this", "that", "are", "was", "have", "has", "not", "but", "all", "can", "will", "的", "了", "在", "是", "和", "有", "为", "等", "年"}
                    doc_keywords -= stop_words

                    for strategy in current_strategies:
                        sid = strategy["id"]
                        content = strategy_contents.get(sid, {})
                        for section in content.get("sections", []):
                            sec_title_lower = section.get("title", "").lower()
                            sec_kw_score = 0
                            for kw in doc_keywords:
                                if len(kw) >= 2 and kw in sec_title_lower:
                                    sec_kw_score += len(kw) * 5
                            if sec_kw_score == 0:
                                continue
                            for item in section.get("items", []):
                                item_text = (item.get("name", "") + " " + item.get("desc", "")).lower()
                                item_score = sec_kw_score
                                for kw in doc_keywords:
                                    if len(kw) >= 2 and kw in item_text:
                                        item_score += len(kw) * 3
                                data_bonus = {"运营数据": 50, "运营策略": 45, "数据与策略": 45, "转化率": 40, "参考数据": 40, "数据统计": 40, "📊": 35, "统计": 30, "概述": 30, "数据": 20}
                                for bkw, bval in data_bonus.items():
                                    if bkw in item_text:
                                        item_score += bval
                                if re.search(r'\d{1,2}月\d{1,2}日', item.get("name", "")):
                                    item_score -= 30
                                if global_best is None or item_score > global_best["score"]:
                                    global_best = {"strategy_id": sid, "strategy": strategy, "section": section, "item": item, "score": item_score}

                if not global_best:
                    # 兜底：找平台策略中"参考数据/数据统计/运营数据"类item，实在没有就用第一个item
                    fallback_bonus_kws = ["参考数据", "数据统计", "运营数据", "数据与策略", "📊", "数据", "概述", "背景"]
                    for strategy in current_strategies:
                        sid = strategy["id"]
                        content_s = strategy_contents.get(sid, {})
                        for section in content_s.get("sections", []):
                            for item in section.get("items", []):
                                item_text = (item.get("name", "") + " " + item.get("desc", "")).lower()
                                for bkw in fallback_bonus_kws:
                                    if bkw in item_text:
                                        global_best = {"strategy_id": sid, "strategy": strategy, "section": section, "item": item, "score": 0}
                                        break
                            if global_best:
                                break
                        if global_best:
                            break
                    # 还是没有就取第一个策略的第一个item
                    if not global_best and current_strategies:
                        first_strategy = current_strategies[0]
                        sid = first_strategy["id"]
                        content_s = strategy_contents.get(sid, {})
                        sections = content_s.get("sections", [])
                        if sections and sections[0].get("items"):
                            global_best = {"strategy_id": sid, "strategy": first_strategy, "section": sections[0], "item": sections[0]["items"][0], "score": 0}
                    if not global_best:
                        continue

                doc_text_stripped = doc_text.strip()
                if len(doc_text_stripped) < 2000:
                    report_summary = doc_text_stripped
                else:
                    parsed = parse_report_content(doc_text, doc_filename)
                    if parsed.get("parsed_items"):
                        report_summary = "\n".join(parsed["parsed_items"])
                    else:
                        report_summary = doc_text_stripped[:500]

                if not report_summary:
                    continue

                best_item = global_best["item"]
                doc_id = doc.get("id", "")
                report_entry = {
                    "game_name": game_name,
                    "data": report_summary,
                    "date": datetime.now().strftime('%Y-%m-%d'),
                    "doc_id": str(doc_id)
                }

                if "report_data" not in best_item:
                    best_item["report_data"] = []
                existing_doc_ids = [rd.get("doc_id") for rd in best_item["report_data"]]
                if str(doc_id) not in existing_doc_ids:
                    best_item["report_data"].append(report_entry)
                    target_sid = global_best["strategy_id"]
                    report_assignments[target_sid] = True

        # 更新策略到数据库
        updated_categories = []
        for strategy in current_strategies:
            sid = strategy["id"]
            s_category = strategy.get("category", "")
            content_modified = sid in report_assignments

            if content_modified:
                content = strategy.get("content", {})
                new_version = strategy.get("version", 1) + 1
                wxcloud_db.update_strategy_content(sid, content, new_version)
                updated_categories.append(s_category)

        # 标记普通文档和报告文档为已应用
        for doc in normal_docs:
            wxcloud_db.update_doc_status(doc["id"], "applied")
        for doc in report_docs:
            wxcloud_db.update_doc_status(doc["id"], "applied")

        doc_type_summary = []
        if normal_docs:
            doc_type_summary.append(f"{len(normal_docs)}个普通文档")
        if report_docs:
            doc_type_summary.append(f"{len(report_docs)}个内部项目报告")

        wxcloud_db.add_update_log(
            platform_id, "批量更新策略",
            f"基于{'、'.join(doc_type_summary)}更新了{len(updated_categories)}个策略分类")

        return {
            "code": 0,
            "msg": f"策略更新成功！基于{'、'.join(doc_type_summary)}更新了{len(updated_categories)}个策略分类",
            "data": {
                "docs_processed": len(pending_docs),
                "categories_updated": updated_categories,
                "report_docs": len(report_docs)
            }
        }
    except Exception as e:
        import traceback
        traceback.print_exc()
        return {"code": -1, "msg": f"更新失败: {str(e)}"}


@app.get("/api/logs")
async def get_logs(limit: int = 20):
    """获取操作日志"""
    try:
        logs = wxcloud_db.get_logs(limit=limit)
        return {"code": 0, "data": logs}
    except Exception as e:
        return {"code": -1, "msg": f"加载日志失败: {str(e)}"}


@app.post("/api/documents/upload-url")
async def upload_url_document(data: UrlUpload):
    """通过网址抓取网页内容并保存为文档"""
    url = data.url.strip()
    if not url:
        raise HTTPException(status_code=400, detail="请输入有效的网址")
    parsed = urlparse(url)
    if not parsed.scheme or not parsed.netloc:
        raise HTTPException(status_code=400, detail="无效的URL格式")
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8"
        }
        resp = requests.get(url, headers=headers, timeout=15, allow_redirects=True)
        resp.raise_for_status()
        if resp.encoding and resp.encoding.lower() != 'utf-8':
            resp.encoding = resp.apparent_encoding
        html_content = resp.text
        title_match = re.search(r'<title[^>]*>(.*?)</title>', html_content, re.IGNORECASE | re.DOTALL)
        page_title = title_match.group(1).strip() if title_match else parsed.netloc
        text = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r'<!--.*?-->', '', text, flags=re.DOTALL)
        text = text.replace('&nbsp;', ' ').replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>').replace('&quot;', '"')
        text = re.sub(r'<[^>]+>', '\n', text)
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = re.sub(r'[ \t]+', ' ', text)
        text = text.strip()
        if len(text) > 50000:
            text = text[:50000] + "\n\n...(内容已截断)"
        doc_title = data.title.strip() if data.title and data.title.strip() else page_title
        full_content = f"来源网址: {url}\n页面标题: {page_title}\n抓取时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n{'='*50}\n\n{text}"
    except requests.exceptions.Timeout:
        raise HTTPException(status_code=408, detail="网页抓取超时")
    except requests.exceptions.ConnectionError:
        raise HTTPException(status_code=502, detail="无法连接到目标网站")
    except requests.exceptions.HTTPError as e:
        raise HTTPException(status_code=502, detail=f"网页请求失败: HTTP {e.response.status_code}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"抓取失败: {str(e)}")

    try:
        result = wxcloud_db.add_document(
            platform_id=data.platform_id,
            filename=doc_title,
            content=full_content,
            doc_type=data.doc_type,
            game_name=data.game_name,
            status="pending"
        )
        doc_id = result.get("doc_id", 0)
        wxcloud_db.add_update_log(
            data.platform_id, f"网址抓取: {doc_title}",
            f"来源: {url}, 文档类型: {data.doc_type}")
        return {
            "code": 0,
            "data": {"id": doc_id, "filename": doc_title, "url": url, "content_length": len(text), "page_title": page_title},
            "msg": f"抓取成功！已保存「{doc_title}」（{len(text)}字符）"
        }
    except Exception as e:
        return {"code": -1, "msg": f"保存失败: {str(e)}"}


@app.get("/api/debug/report-diagnosis")
async def debug_report_diagnosis():
    """诊断接口：查看数据库中报告文档和策略中报告数据的完整状态"""
    try:
        # 1. 查询 uploaded_docs 表中所有 report 类型文档
        with wxcloud_db.get_db() as db:
            cursor = db.cursor()
            cursor.execute("SELECT id, platform_id, filename, doc_type, game_name, status, created_at FROM uploaded_docs WHERE doc_type='report' ORDER BY created_at DESC")
            report_docs = cursor.fetchall()
            for d in report_docs:
                d["created_at"] = wxcloud_db._dt_to_str(d.get("created_at"))
            
            # 2. 查询各类型文档数量
            cursor.execute("SELECT doc_type, COUNT(*) as cnt FROM uploaded_docs GROUP BY doc_type")
            doc_type_counts = cursor.fetchall()
            
            # 3. 查询所有文档总数
            cursor.execute("SELECT COUNT(*) as total FROM uploaded_docs")
            total_docs = cursor.fetchone()["total"]
            
            cursor.close()
        
        # 4. 检查策略中嵌入的 report_data
        strategies = wxcloud_db.get_all_strategies()
        strategies_with_reports = []
        for s in strategies:
            content = s.get("content", {})
            for section in content.get("sections", []):
                for item in section.get("items", []):
                    rd_list = item.get("report_data", [])
                    if rd_list:
                        strategies_with_reports.append({
                            "strategy_id": s["id"],
                            "platform_name": s.get("platform_name", ""),
                            "category": s.get("category", ""),
                            "section": section.get("title", ""),
                            "item": item.get("name", ""),
                            "report_data_count": len(rd_list),
                            "report_entries": rd_list
                        })
        
        return {
            "code": 0,
            "data": {
                "total_docs_in_db": total_docs,
                "doc_type_distribution": doc_type_counts,
                "report_docs_in_uploaded_docs": len(report_docs),
                "report_docs": report_docs,
                "strategies_with_embedded_reports": len(strategies_with_reports),
                "embedded_report_details": strategies_with_reports
            }
        }
    except Exception as e:
        import traceback
        return {"code": -1, "msg": f"诊断失败: {str(e)}", "trace": traceback.format_exc()}


@app.post("/api/debug/reset-report-status")
async def reset_report_status():
    """将所有内部项目报告重置为pending状态"""
    try:
        where = '{doc_type: "report"}'
        wxcloud_db.db_update("uploaded_docs", where, {"status": "pending", "updated_at": wxcloud_db._now_str()})
        return {"code": 0, "msg": "已将内部项目报告重置为待处理状态"}
    except Exception as e:
        return {"code": -1, "msg": f"重置失败: {str(e)}"}


class ReportDataEdit(BaseModel):
    strategy_id: int
    section_index: int
    item_index: int
    report_data_index: int
    game_name: Optional[str] = None
    data: Optional[str] = None
    date: Optional[str] = None


class ReportDataDelete(BaseModel):
    strategy_id: int
    section_index: int
    item_index: int
    report_data_index: int


@app.post("/api/strategies/edit-report-data")
async def edit_report_data(edit: ReportDataEdit):
    """编辑策略中的内部项目数据条目"""
    try:
        strategy = wxcloud_db.get_strategy_by_id(edit.strategy_id)
        if not strategy:
            return {"code": -1, "msg": "策略不存在"}

        content = strategy.get("content", {})
        if not content:
            content = {"sections": []}

        sections = content.get("sections", [])
        if edit.section_index < 0 or edit.section_index >= len(sections):
            return {"code": -1, "msg": "section索引无效"}
        items = sections[edit.section_index].get("items", [])
        if edit.item_index < 0 or edit.item_index >= len(items):
            return {"code": -1, "msg": "item索引无效"}
        report_data = items[edit.item_index].get("report_data", [])
        if edit.report_data_index < 0 or edit.report_data_index >= len(report_data):
            return {"code": -1, "msg": "report_data索引无效"}

        rd = report_data[edit.report_data_index]
        if edit.game_name is not None:
            rd["game_name"] = edit.game_name
        if edit.data is not None:
            rd["data"] = edit.data
        if edit.date is not None:
            rd["date"] = edit.date

        new_version = strategy.get("version", 1) + 1
        wxcloud_db.update_strategy_content(edit.strategy_id, content, new_version)
        wxcloud_db.add_update_log(
            strategy.get("platform_id"), "编辑内部项目数据",
            f"策略ID: {edit.strategy_id}")

        return {"code": 0, "msg": "更新成功", "data": rd}
    except Exception as e:
        return {"code": -1, "msg": f"更新失败: {str(e)}"}


@app.post("/api/strategies/delete-report-data")
async def delete_report_data(delete: ReportDataDelete):
    """删除策略中的内部项目数据条目"""
    try:
        strategy = wxcloud_db.get_strategy_by_id(delete.strategy_id)
        if not strategy:
            return {"code": -1, "msg": "策略不存在"}

        content = strategy.get("content", {})
        if not content:
            content = {"sections": []}

        sections = content.get("sections", [])
        if delete.section_index < 0 or delete.section_index >= len(sections):
            return {"code": -1, "msg": "section索引无效"}
        items = sections[delete.section_index].get("items", [])
        if delete.item_index < 0 or delete.item_index >= len(items):
            return {"code": -1, "msg": "item索引无效"}
        report_data = items[delete.item_index].get("report_data", [])
        if delete.report_data_index < 0 or delete.report_data_index >= len(report_data):
            return {"code": -1, "msg": "report_data索引无效"}

        report_data.pop(delete.report_data_index)

        new_version = strategy.get("version", 1) + 1
        wxcloud_db.update_strategy_content(delete.strategy_id, content, new_version)
        wxcloud_db.add_update_log(
            strategy.get("platform_id"), "删除内部项目数据",
            f"策略ID: {delete.strategy_id}")

        return {"code": 0, "msg": "删除成功"}
    except Exception as e:
        return {"code": -1, "msg": f"删除失败: {str(e)}"}


class ReportPasswordVerify(BaseModel):
    password: str

REPORT_VIEW_PASSWORD = os.environ.get("REPORT_VIEW_PASSWORD", "YOUR_INTERNAL_PASSWORD")


@app.post("/api/verify-report-password")
async def verify_report_password(data: ReportPasswordVerify):
    """验证内部项目数据查看密码"""
    if not data.password or not data.password.strip():
        return {"code": -1, "msg": "请输入密码"}
    if data.password.strip() == REPORT_VIEW_PASSWORD:
        return {"code": 0, "msg": "验证成功"}
    else:
        return {"code": -1, "msg": "密码错误，请重试"}


class AskQuestion(BaseModel):
    question: str


@app.post("/api/ask")
async def ask_question(data: AskQuestion):
    """智能问答"""
    question = data.question.strip()
    if not question:
        return {"code": -1, "msg": "请输入问题"}

    try:
        platforms = wxcloud_db.get_all_platforms()
        strategies = wxcloud_db.get_all_strategies()

        # 获取文档列表（用于搜索）
        docs_result = wxcloud_db.db_query("uploaded_docs", order_by="-created_at", limit=50)
        docs = docs_result.get("data", [])

        q_lower = question.lower()
        category_map = {"promotion": "推广资源", "operation": "运营功能", "technology": "技术合作"}

        platform_keywords = {
            "steam": ["steam", "v社", "valve", "steamworks", "steam deck"],
            "epic": ["epic", "epic games", "虚幻", "eos", "unreal"],
            "xbox": ["xbox", "微软", "microsoft", "game pass", "gamepass", "xgp", "showcase", "developer direct"],
            "playstation": ["playstation", "ps", "ps5", "ps4", "索尼", "sony", "psn", "ps plus"]
        }

        category_keywords = {
            "promotion": ["推广", "促销", "折扣", "打折", "优惠", "活动", "sale", "特卖", "降价", "资源", "曝光"],
            "operation": ["运营", "功能", "上架", "发布", "审核", "分成", "收入", "key", "密钥", "愿望单"],
            "technology": ["技术", "sdk", "api", "适配", "手柄", "触觉", "deck", "verified", "跨平台"]
        }

        topic_keywords = {
            "折扣力度": ["折扣力度", "打几折", "多少折", "折扣幅度", "定价"],
            "时间安排": ["什么时候", "何时", "时间", "日期", "日历", "开始", "结束"],
            "入选条件": ["怎么参加", "如何参加", "条件", "要求", "资格", "申请"],
            "分成收入": ["分成", "收入", "佣金", "比例", "revenue"]
        }

        matched_platforms = []
        for pname, kws in platform_keywords.items():
            for kw in kws:
                if kw in q_lower:
                    matched_platforms.append(pname)
                    break

        matched_categories = []
        for cat, kws in category_keywords.items():
            for kw in kws:
                if kw in q_lower:
                    matched_categories.append(cat)
                    break

        all_q_words = re.findall(r'[\u4e00-\u9fff]+|[a-zA-Z]+', q_lower)
        platform_only_words = set()
        for pname, kws in platform_keywords.items():
            for kw in kws:
                platform_only_words.add(kw)
        core_words = [w for w in all_q_words if len(w) >= 2 and w.lower() not in platform_only_words]
        if not core_words:
            core_words = [w for w in all_q_words if len(w) >= 2]

        if not matched_platforms:
            target_platforms = [p["name"].lower().replace(" ", "") for p in platforms]
        else:
            target_platforms = matched_platforms

        matched_strategy_items = []
        for strategy in strategies:
            p_name = (strategy.get("platform_name") or "").lower().replace(" ", "")
            s_cat = strategy.get("category", "")

            platform_match = False
            for tp in target_platforms:
                if tp in p_name or p_name in tp:
                    platform_match = True
                    break
            if not platform_match and matched_platforms:
                continue
            if matched_categories and s_cat not in matched_categories:
                continue

            content = strategy.get("content", {})
            for section in content.get("sections", []):
                section_title = section.get("title", "")
                for item in section.get("items", []):
                    item_name = item.get("name", "")
                    item_desc = item.get("desc", "")
                    combined_text = f"{section_title} {item_name} {item_desc}".lower()

                    score = 0
                    core_hits = 0
                    for w in core_words:
                        if w in combined_text:
                            score += 10
                            core_hits += 1

                    if len(core_words) >= 2:
                        phrase_nospace = "".join(core_words)
                        if phrase_nospace in combined_text:
                            score += 30
                        elif core_hits == len(core_words):
                            score += 15

                    if core_hits == 0 and core_words:
                        continue
                    if s_cat in matched_categories:
                        score += 2

                    if score > 0:
                        matched_strategy_items.append({
                            "score": score, "core_hits": core_hits,
                            "platform": strategy.get("platform_name", ""),
                            "platform_icon": strategy.get("platform_icon", "🎮"),
                            "category": category_map.get(s_cat, s_cat),
                            "section": section_title,
                            "name": item_name, "desc": item_desc,
                            "source": item.get("source", ""),
                            "game_name": item.get("game_name", ""),
                            "report_data": item.get("report_data", [])
                        })

        matched_strategy_items.sort(key=lambda x: x["score"], reverse=True)

        if matched_strategy_items:
            max_score = matched_strategy_items[0]["score"]
            threshold = max(max_score * 0.4, 5)
            top_results = [item for item in matched_strategy_items if item["score"] >= threshold][:10]
        else:
            top_results = []

        matched_docs = []
        plat_map = {p["id"]: p["name"] for p in platforms}
        for doc in docs:
            doc_text = f"{doc.get('filename', '')} {(doc.get('content', '') or '')[:500]}".lower()
            doc_score = 0
            doc_core_hits = 0
            for w in core_words:
                if w in doc_text:
                    doc_score += 3
                    doc_core_hits += 1
            if doc_core_hits == 0 and core_words:
                continue
            if doc_score > 0:
                matched_docs.append({
                    "score": doc_score,
                    "filename": doc.get("filename", ""),
                    "platform_name": plat_map.get(doc.get("platform_id"), "通用"),
                    "doc_type": doc.get("doc_type", ""),
                    "status": doc.get("status", ""),
                    "created_at": doc.get("created_at", "")
                })
        matched_docs.sort(key=lambda x: x["score"], reverse=True)
        top_docs = matched_docs[:5]

        if not top_results and not top_docs:
            summary_parts = []
            for p in platforms:
                p_strats = [s for s in strategies if s.get("platform_id") == p["id"]]
                cats = [category_map.get(s["category"], s["category"]) for s in p_strats]
                summary_parts.append(f"{p.get('icon', '🎮')} **{p['name']}**：共{len(p_strats)}项策略（{'、'.join(cats)}）")
            answer = (
                f"抱歉，未找到与「{question}」直接相关的策略内容。\n\n"
                "以下是当前系统中的全部平台策略概览：\n\n" +
                "\n".join(summary_parts) +
                "\n\n💡 **提问建议**：\n"
                "• 指定平台名，如「Steam春促什么时候？」\n"
                "• 询问具体内容，如「PS5游戏如何参加黑五折扣？」"
            )
            return {"code": 0, "data": {"answer": answer, "sources": [], "docs": []}}

        platform_groups = {}
        for item in top_results:
            p = item["platform"]
            if p not in platform_groups:
                platform_groups[p] = {"icon": item["platform_icon"], "items": []}
            platform_groups[p]["items"].append(item)

        answer_lines = [f"关于「{question}」，以下是为您检索到的相关策略信息：\n"]
        for p_name, group in platform_groups.items():
            answer_lines.append(f"\n{group['icon']} **{p_name}**\n")
            section_groups = {}
            for item in group["items"]:
                sec = item["section"]
                if sec not in section_groups:
                    section_groups[sec] = []
                section_groups[sec].append(item)
            for sec_title, items in section_groups.items():
                answer_lines.append(f"📌 {sec_title}")
                for item in items:
                    source_tag = ""
                    if item.get("source") == "report":
                        game = item.get("game_name", "")
                        source_tag = f" 🏷️[内部项目报告 · {game}]" if game else " 🏷️[内部项目报告]"
                    answer_lines.append(f"• **{item['name']}**：{item['desc']}{source_tag}")
                    for rd in item.get("report_data", []):
                        answer_lines.append(f"  📊 **{rd.get('game_name', '未知项目')}** ({rd.get('date', '')})：{rd.get('data', '')}")
                answer_lines.append("")

        if top_docs:
            answer_lines.append("\n📄 **相关参考文档**：")
            for doc in top_docs:
                status_text = "待处理" if doc["status"] == "pending" else "已应用"
                answer_lines.append(f"• [{status_text}] {doc['filename']}（{doc['platform_name']}，{doc['created_at']}）")

        answer = "\n".join(answer_lines)
        sources = []
        seen = set()
        for item in top_results:
            key = f"{item['platform']}-{item['category']}"
            if key not in seen:
                seen.add(key)
                sources.append({"platform": item["platform"], "category": item["category"], "icon": item["platform_icon"]})

        return {
            "code": 0,
            "data": {"answer": answer, "sources": sources, "docs": top_docs, "total_matches": len(matched_strategy_items)}
        }
    except Exception as e:
        import traceback
        traceback.print_exc()
        return {"code": -1, "msg": f"查询失败: {str(e)}"}


## ==================== 静态文件 & 重定向 ====================
from fastapi.staticfiles import StaticFiles
from fastapi.responses import RedirectResponse as _Redirect
from fastapi.responses import FileResponse

# 确保 static 目录存在（防止 Docker 构建遗漏导致启动崩溃）
os.makedirs("static", exist_ok=True)

@app.get("/")
async def root():
    """根路径重定向到平台页面"""
    return _Redirect("/platform/static/index.html")

@app.get("/platform/static/{path:path}")
async def serve_platform_static(path: str):
    """兼容nginx路径，直接提供静态文件"""
    file_path = os.path.join("static", path)
    if os.path.isfile(file_path):
        # 根据扩展名设置正确的MIME类型
        ext = path.rsplit('.', 1)[-1].lower() if '.' in path else ''
        mime_map = {
            'html': 'text/html', 'css': 'text/css', 'js': 'application/javascript',
            'json': 'application/json', 'png': 'image/png', 'jpg': 'image/jpeg',
            'svg': 'image/svg+xml', 'ico': 'image/x-icon', 'woff2': 'font/woff2',
        }
        return FileResponse(file_path, media_type=mime_map.get(ext, 'application/octet-stream'))
    return JSONResponse(status_code=404, content={"detail": "Not found"})

# 必须放在最后，避免路径冲突
app.mount("/static", StaticFiles(directory="static", html=True), name="static")


if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)